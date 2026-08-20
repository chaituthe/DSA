import os
import mss
import pandas as pd
from datetime import datetime
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
from tkinter import messagebox

capture_log = pd.DataFrame(columns=[
    "monitor",
    "file_name",
    "limits",
    "source",
    "status",
    "test_value",
    "capture_time",
    "output_path"
])

def capture_region(monitor_index, width, height, output_name="screenshot.png", output_dir=".", right_margin=0, bottom_margin=0):
    os.makedirs(output_dir, exist_ok=True)

    with mss.mss() as sct:
        if monitor_index >= len(sct.monitors):
            print(f"Error: Monitor {monitor_index} not found. You have {len(sct.monitors) - 1} monitors connected.")
            return None

        monitor = sct.monitors[monitor_index]

        left = monitor["left"] + monitor["width"] - width - right_margin
        top = monitor["top"] + monitor["height"] - height - bottom_margin

        capture_box = {
            "top": top,
            "left": left,
            "width": width,
            "height": height
        }

        screenshot = sct.grab(capture_box)
        img = Image.frombytes("RGBA", screenshot.size, screenshot.bgra, "raw", "BGRA").convert("RGB")

        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        base_name, ext = os.path.splitext(output_name)
        output_path = os.path.join(output_dir, f"{base_name}_{timestamp}{ext}")

        img.save(output_path)
        print(f"Saved screenshot successfully as '{output_path}'")
        return output_path

def append_capture_record(monitor, file_name, limits, source, status, test_value, output_path):
    global capture_log

    row = {
        "monitor": monitor,
        "file_name": file_name,
        "limits": limits,
        "source": source,
        "status": status,
        "test_value": test_value,
        "capture_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "output_path": output_path
    }

    capture_log = pd.concat(
        [capture_log, pd.DataFrame([row])],
        ignore_index=True
    )

    timestamp =""
    #timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    csv_path = os.path.join(r"D:\Learning\Project_One", f"log_{timestamp}.csv")
    os.makedirs(os.path.dirname(csv_path), exist_ok=True)
    capture_log.to_csv(csv_path, index=False)

    print("Saved DataFrame to CSV:", csv_path)


def add_screenshot_to_powerpoint(image_path, pptx_path=r"D:\Learning\Project_One\power.pptx"):
    try:
        if os.path.exists(pptx_path):
            presentation = Presentation(pptx_path)
        else:
            presentation = Presentation()

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        picture = slide.shapes.add_picture(image_path, Inches(0.3), Inches(0.3))

        max_width = presentation.slide_width - Inches(0.6)
        max_height = presentation.slide_height - Inches(0.6)

        if picture.width > max_width or picture.height > max_height:
            ratio = min(max_width / picture.width, max_height / picture.height)
            picture.width = int(picture.width * ratio)
            picture.height = int(picture.height * ratio)

        presentation.save(pptx_path)
        print(f"Inserted screenshot into PowerPoint: {pptx_path}")
        return pptx_path
    except Exception as e:
        print(f"PowerPoint update failed: {e}")
        raise


def capture_from_gui(selected_monitor_var, selected_file_var, selected_limits_var, selected_source_var, selected_status_var, selected_test_var):
    output_dir = r"D:\Learning\Project_One\screenshots"

    try:
        monitor_index = int(selected_monitor_var.get())
        file_name = selected_file_var.get()
        limits = selected_limits_var.get()
        source = selected_source_var.get()
        status = selected_status_var.get()
        test_value = selected_test_var.get()



        output_path = capture_region(
            monitor_index=monitor_index,
            width=800,
            height=600,
            output_name=file_name,
            output_dir=output_dir,
            right_margin=50,
            bottom_margin=200
        )

        if output_path:
            append_capture_record(
                monitor=monitor_index,
                file_name=file_name,
                limits=limits,
                source=source,
                status=status,
                test_value=test_value,
                output_path=output_path
            )

            add_screenshot_to_powerpoint(output_path)

            print(capture_log.tail(1).to_dict(orient="records"))
            print(capture_log)

        print(f"Selected monitor: {monitor_index}")
        print(f"Selected file: {file_name}")
        print(f"Selected limits: {limits}")
        print(f"Selected source: {source}")
        print(f"Selected status: {status}")
        print(f"Test box value: {test_value}")
    except Exception as e:
        messagebox.showerror("Error", str(e))

def create_selection_group(parent, title, options, selected_var, prefix=""):
    group = tk.LabelFrame(parent, text=title, padx=10, pady=8, bd=1, relief="solid")
    group.pack(fill="both", expand=True, padx=10, pady=6)

    for option in options:
        value = str(option)
        label = f"{prefix} {value}" if prefix else value

        tk.Radiobutton(
            group,
            text=label,
            variable=selected_var,
            value=value
        ).pack(anchor="w")

    return group

def build_gui():
    root = tk.Tk()
    root.title("Screen Capture")
    root.geometry("500x650")
    root.resizable(False, False)

    with mss.mss() as sct:
        monitor_choices = [str(i) for i in range(1, len(sct.monitors))]

    if not monitor_choices:
        monitor_choices = ["1"]

    selected_monitor = tk.StringVar(value=monitor_choices[0])
    selected_file = tk.StringVar(value="22.png")
    selected_limits = tk.StringVar(value="DALP")
    selected_source = tk.StringVar(value="DALP")
    selected_status = tk.StringVar(value="active")
    selected_test = tk.StringVar(value="Enter test value")

    file_options = [
        "22.png",
        "33.png",
        "44.png",
        "55.png"
    ]

    limits_options = [
        "s22",
        "s33",
        "s44",
        "s55",

    ]

    source_options = [
        "c22",
        "c33",
        "c44",
        "c55"
    ]
    

    status_options = [
        "d22",
        "d33",
        "d24",
        "d55"
    ]

    main_frame = tk.Frame(root)
    main_frame.pack(fill="both", expand=True, padx=10, pady=10)

    left_column = tk.Frame(main_frame)
    left_column.grid(row=0, column=0, sticky="nsew", padx=(0, 10))

    right_column = tk.Frame(main_frame)
    right_column.grid(row=0, column=1, sticky="nsew")

    main_frame.grid_columnconfigure(0, weight=1)
    main_frame.grid_columnconfigure(1, weight=1)

    create_selection_group(left_column, "Monitor", monitor_choices, selected_monitor, "Monitor")
    create_selection_group(left_column, "File name", file_options, selected_file)

    create_selection_group(right_column, "Limits ", limits_options, selected_limits)
    create_selection_group(right_column, "Source", source_options, selected_source)
    create_selection_group(right_column, "Status", status_options, selected_status)

    test_frame = tk.Frame(root)
    test_frame.pack(fill="x", padx=10, pady=(0, 10))

    tk.Label(test_frame, text="Test Box:").pack(side="left", padx=(0, 8))
    tk.Entry(test_frame, textvariable=selected_test, width=30).pack(side="left", fill="x", expand=True)

    button = tk.Button(
        root,
        text="Capture Image",
        width=20,
        height=2,
        command=lambda: capture_from_gui(
            selected_monitor,
            selected_file,
            selected_limits,
            selected_source,
            selected_status,
            selected_test
        )
    )
    button.pack(pady=15)

    root.mainloop()

if __name__ == "__main__":
    build_gui()
